"""
Generates three sample investor pitch deck PPTX files for PitchFlip demos.
Output: frontend/public/samples/{aurica,lexara,carbonthread}.pptx

Run from the repo root:
    cd pitch-product && python sample-decks/generate_pptx.py
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_DIR = Path(__file__).parent.parent / "frontend" / "public" / "samples"
OUT_DIR.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_slide(prs: Presentation, layout_index: int = 6) -> object:
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def add_text_box(slide, text: str, left: float, top: float, width: float, height: float,
                 font_size: int, bold: bool = False, color: str = "#FFFFFF",
                 align=PP_ALIGN.LEFT, italic: bool = False) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)


def add_background(slide, hex_color: str) -> None:
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def add_content_slide(prs, bg: str, accent: str, title: str, bullets: list[str]) -> None:
    slide = add_slide(prs)
    add_background(slide, bg)
    add_text_box(slide, title, 0.5, 0.3, 9.0, 0.8, font_size=28, bold=True, color=accent)
    # Divider line approximated as a thin rectangle
    from pptx.util import Inches, Pt
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(accent)
    line.line.fill.background()

    bullet_text = "\n".join(f"  {b}" for b in bullets)
    add_text_box(slide, bullet_text, 0.5, 1.4, 9.0, 5.5, font_size=18, color="#E8E8F0")


def add_title_slide(prs, bg: str, accent: str, company: str, tagline: str, subtitle: str) -> None:
    slide = add_slide(prs)
    add_background(slide, bg)
    add_text_box(slide, company, 0.7, 1.8, 8.6, 1.2, font_size=54, bold=True,
                 color="#FFFFFF", align=PP_ALIGN.CENTER)
    add_text_box(slide, tagline, 0.7, 3.2, 8.6, 0.9, font_size=24, bold=False,
                 color=accent, align=PP_ALIGN.CENTER, italic=True)
    add_text_box(slide, subtitle, 0.7, 4.3, 8.6, 0.6, font_size=14,
                 color="#9090A8", align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Deck definitions
# ---------------------------------------------------------------------------

DECKS = [
    {
        "filename": "vantara.pptx",
        "bg": "#1A1000",
        "accent": "#F59E0B",
        "company": "Vantara",
        "tagline": "See supplier risk before it becomes a supply chain crisis.",
        "subtitle": "Series A Pitch  |  2025",
        "slides": [
            ("The Problem", [
                "Supply chain disruptions cost Global 2000 companies an average of $182M per incident",
                "82% of CPOs say they had no warning before their last major supplier failure",
                "Current tools: annual audits, manual credit checks, and spreadsheets built by analysts",
                "By the time a risk is visible — financial distress, geopolitical exposure, quality failure — it's already a crisis",
                "Procurement teams are managing $50M–$10B in supplier spend with effectively zero real-time intelligence",
            ]),
            ("Our Solution", [
                "Vantara monitors your entire supplier base across 150+ risk signals — continuously, not annually",
                "Financial distress signals: payment delays, credit rating changes, court filings, executive departures",
                "Operational signals: production capacity changes, quality incidents, certifications lapsing",
                "Geopolitical signals: sanctions, export controls, conflict proximity, regulatory shifts by country",
                "Risk score per supplier updated daily — with plain-English explanation and recommended action",
            ]),
            ("How It Works", [
                "1. Connect: Upload your supplier list or integrate with SAP Ariba, Coupa, or Oracle Procurement",
                "2. Map: Vantara builds your full supplier network including tier-2 and tier-3 dependencies",
                "3. Monitor: 150+ signals tracked per supplier from 4,000+ data sources, updated daily",
                "4. Alert: Critical risks surface as prioritised alerts with context — not a list of 500 flags",
                "5. Act: Built-in playbooks for each risk type: contingency sourcing, contractual triggers, escalation paths",
            ]),
            ("Traction", [
                "13 enterprise customers across aerospace, defence, automotive, and industrial manufacturing",
                "$3.1M ARR | 155% net revenue retention",
                "Average time from onboarding to first critical supplier alert: 11 days",
                "Customers have collectively avoided an estimated $340M in supply chain disruption costs",
                "Pipeline: $8.4M of qualified enterprise opportunities, predominantly FTSE 100 and Fortune 500",
            ]),
            ("Market Opportunity", [
                "Supply chain risk management market: $6.8B in 2024, projected $18.4B by 2030 (CAGR 18%)",
                "EU Corporate Sustainability Due Diligence Directive (CSDDD) mandates supplier risk monitoring for 50,000 companies",
                "US CHIPS Act and defence procurement regulations requiring verified supplier resilience",
                "Average enterprise with 500+ suppliers spends $2M+/year on manual supplier audits today",
                "Every large manufacturer, defence contractor, and critical infrastructure company is our prospect",
            ]),
            ("Business Model", [
                "Per-supplier monitoring: $28/active supplier/month",
                "Platform fee: $4,000/month (includes integrations, alerting, playbooks)",
                "Average enterprise: 350 active suppliers = $9,800/month + platform = $165,000 ARR",
                "Tier-2/3 network mapping: additional $2,500/month per extended tier",
                "Average contract length: 3 years | Gross margin: 81%",
            ]),
            ("Why We Win", [
                "Dun & Bradstreet / Moody's: financial data only. No operational or geopolitical signals. No action layer.",
                "Riskmethods / Resilinc: legacy tools built for procurement analysts, not CPOs. High implementation cost.",
                "Manual audits: annual snapshots. $5,000–$20,000 per supplier. Zero continuous monitoring.",
                "Vantara is the only product combining financial, operational, and geopolitical signals in a single daily-updated risk score",
                "Proprietary signal network trained on 6,000+ historical supplier failure events — accuracy improves with scale",
            ]),
            ("Team", [
                "CEO — Clara Voss: Former Head of Supply Chain Strategy at Airbus. Led a €2.4B supplier risk programme across 1,800 suppliers in 40 countries. Previously at McKinsey Operations Practice.",
                "CTO — Dr. Rahul Mehta: PhD in network risk modelling from MIT. Former Principal Data Scientist at Palantir, where he built supply chain intelligence for three NATO defence programmes.",
                "CPO — James Okafor: Former VP Product at Coupa Software. Launched their supplier risk module. Previously Head of Product at Tradeshift.",
                "Advisory board: Former CPO of BAE Systems, Chief Risk Officer of Maersk, ex-Director of the UK National Cyber Security Centre.",
            ]),
            ("The Ask", [
                "Raising $12M Series A",
                "Use of funds: 50% product (tier-2/3 mapping, geopolitical signal layer) | 35% enterprise sales (UK, EU, US) | 15% operations",
                "Two investors committed at $7M combined. Seeking one strategic co-investor from aerospace, defence, or industrial manufacturing.",
                "Regulatory tailwinds: CSDDD enforcement and US supply chain resilience mandates are creating non-discretionary spend",
                "Milestones: $11M ARR within 18 months, SAP Ariba and Coupa marketplace listings live",
            ]),
        ],
    },
    {
        "filename": "lexara.pptx",
        "bg": "#0A1628",
        "accent": "#38BDF8",
        "company": "Lexara",
        "tagline": "Contract risk, surfaced in seconds.",
        "subtitle": "Series A Pitch  |  2025",
        "slides": [
            ("The Problem", [
                "Enterprise procurement teams review 500–2,000 contracts per year",
                "Standard review with outside counsel: 3–5 days and $2,000–$8,000 per contract",
                "1 in 4 contracts contains a material risk clause missed during review",
                "Legal teams are bottlenecks — business deals stall waiting for sign-off",
                "The result: slower deal cycles, higher legal spend, and unquantified exposure",
            ]),
            ("Our Solution", [
                "Lexara analyses every clause in a contract in under 4 minutes",
                "Trained on 6 million contracts across 14 industries — understands context, not just keywords",
                "Surfaces unfavourable terms, missing protections, and benchmarks against market standards",
                "Integrates with DocuSign, SharePoint, iManage, and all major CLM platforms",
                "Legal teams review only the 20% of contracts that need human attention",
            ]),
            ("How It Works", [
                "1. Upload or connect your contract repository (DocuSign, SharePoint, email)",
                "2. Lexara extracts and classifies every clause: liability, IP, termination, payment, indemnity",
                "3. Risk score assigned 1–10 with plain-English explanations for each flag",
                "4. Benchmark: each clause compared against 6M contracts in your industry",
                "5. Route: only contracts above your risk threshold go to a human reviewer",
            ]),
            ("Traction", [
                "9 enterprise customers including two Magic Circle law firms and three FTSE 250 procurement teams",
                "$2.1M ARR | 140% net revenue retention",
                "Average time-to-review reduced from 4.2 days to 18 minutes",
                "Average legal spend reduction per customer: $840,000 annually",
                "Pipeline: $6.2M of qualified enterprise deals in active procurement cycles",
            ]),
            ("Market Opportunity", [
                "Global legal tech market: $27B in 2024, projected $50B by 2029",
                "Contract lifecycle management segment: $3.2B, growing at 14% CAGR",
                "Enterprise procurement legal spend: $185B annually — almost entirely manual today",
                "CSRD and supply chain due diligence regulation creating new contract review mandates",
                "Every Fortune 1000 company is our prospect",
            ]),
            ("Business Model", [
                "Per-seat SaaS: $600/user/month (legal and procurement users)",
                "Volume pricing: $1,200/month per 1,000 contracts processed",
                "Implementation fee: $15,000 one-time for enterprise deployments",
                "Average enterprise contract value: $240,000 ARR",
                "Gross margin: 82%",
            ]),
            ("Why We Win", [
                "Kira Systems / Luminance: document review tools. Require months of specialist training. No risk scoring.",
                "DocuSign CLM: workflow management only. No risk intelligence layer.",
                "Outside counsel: 10–50x more expensive. Creates dependency. Not scalable.",
                "Lexara is the only product combining clause extraction, market benchmarking, and plain-English risk explanation",
                "First product built by a practising partner (not technologists guessing at legal workflow)",
            ]),
            ("Team", [
                "CEO — Sophia Andersen: Former partner at Clifford Chance LLP. Led the firm's legal technology practice for 6 years. Advised 40 FTSE 100 companies on contract risk.",
                "CTO — Dr. Kevin Yip: PhD in NLP from Stanford. Former research lead at LexisNexis AI. 11 patents in legal language modelling.",
                "CCO — Marcus Webb: Former Head of Legal Operations at Vodafone. Reduced their external legal spend by 64% using automation. Scaled in-house team from 40 to 12.",
                "Advisory board: GC of Rolls-Royce, former Chief Innovation Officer of Allen & Overy.",
            ]),
            ("The Ask", [
                "Raising $11M Series A",
                "Use of funds: 55% product and engineering | 30% enterprise sales and legal partnerships | 15% operations",
                "Lead term sheet received. Seeking one or two co-investors with enterprise SaaS or legal sector expertise.",
                "Milestones: $10M ARR within 18 months, 25 enterprise customers",
                "Regulatory tailwinds (EU AI Act, CSRD) accelerating inbound demand — pipeline doubled in Q1 2025",
            ]),
        ],
    },
    {
        "filename": "carbonthread.pptx",
        "bg": "#0A1F0F",
        "accent": "#4ADE80",
        "company": "Carbonthread",
        "tagline": "Your supply chain's carbon footprint, finally visible.",
        "subtitle": "Series A Pitch  |  2025",
        "slides": [
            ("The Problem", [
                "70% of a company's carbon footprint lives in Scope 3 — emissions from suppliers they don't control",
                "Current approach: annual surveys sent to suppliers. Average response rate: 38%.",
                "Data is months out of date, unauditable, and fails regulatory scrutiny",
                "CSRD (EU) and SEC climate disclosure rules require verified Scope 3 data — starting now",
                "Companies face fines, ESG rating downgrades, and investor withdrawal with no reliable data to show",
            ]),
            ("Our Solution", [
                "Carbonthread connects directly to ERP, procurement, and logistics systems",
                "Maps every supplier transaction to a verified carbon emissions factor — automatically, in real time",
                "Aggregates Scope 1, 2, and 3 into a single auditable dashboard",
                "Supplier engagement module: automated data requests, peer benchmarking, gap analysis",
                "One-click export for CSRD, GRI, TCFD, and CDP — audit-ready on day one",
            ]),
            ("How It Works", [
                "1. Connect: SAP, Oracle, or NetSuite integration — live in 2 weeks, no data migration required",
                "2. Map: Every purchase order matched to one of 80,000+ product and service carbon factors",
                "3. Calculate: Real-time Scope 3 emissions by supplier, category, geography, and product line",
                "4. Engage: Automated supplier data requests for primary data where spend warrants it",
                "5. Report: Regulatory-ready export with methodology footnotes and confidence intervals",
            ]),
            ("Traction", [
                "11 enterprise customers including 3 FTSE 100 companies and 4 mid-market manufacturers",
                "$3.4M ARR | 160% net revenue retention",
                "Average Scope 3 data completeness improved from 18% to 94% within 90 days",
                "Customers have collectively avoided €12M in projected ESG reporting penalties",
                "Expansion pipeline: 7 existing customers upgrading to Supplier Engagement module",
            ]),
            ("Market Opportunity", [
                "ESG reporting software market: $1.6B in 2024, projected $7.5B by 2030 (CAGR 29%)",
                "CSRD brings 50,000 EU companies into mandatory reporting — starting 2025",
                "SEC climate disclosure rules targeting 10,000+ US public companies",
                "Total addressable opportunity in Scope 3 data infrastructure: $18B globally",
                "Regulatory mandates create non-discretionary spend — recession-resilient category",
            ]),
            ("Business Model", [
                "Platform fee: $3,500/month base (includes ERP integration and core dashboard)",
                "Per-supplier pricing: $12/active supplier/month",
                "Average enterprise: 400 active suppliers = $4,800/month + platform = ~$100,000 ARR",
                "Implementation: $20,000 one-time for SAP/Oracle connectors",
                "Gross margin: 79% | Average contract length: 3 years",
            ]),
            ("Why We Win", [
                "Watershed / Persefoni: US-focused, Scope 2 heavy. Weak on Scope 3 supplier data. No ERP integration.",
                "Manual supplier surveys: 38% response rate, unauditable, months out of date.",
                "Big 4 consulting: $500K–$2M per engagement. Annual snapshots. Not continuous.",
                "Carbonthread is the only product that pulls live Scope 3 data from ERP transactions — no surveys required",
                "SAP and Oracle partnership discussions underway — potential to embed as default sustainability module",
            ]),
            ("Team", [
                "CEO — Amara Osei: Former sustainability director at McKinsey. Led the CDP disclosure programme for 40 global clients. Architect of the EU taxonomy mapping methodology now used by three major banks.",
                "CTO — Lars Eriksson: Former engineering lead at SAP's Sustainability division. Designed the carbon accounting layer used by 800 companies. Co-author of ISO 14064-3 guidance.",
                "CPO — Nina Petrov: Former lead product manager at Salesforce Sustainability Cloud. Launched their CSRD module. Previously Head of Product at Persefoni.",
                "Advisory board: Former DG of DG CLIMA (EU), Chief Sustainability Officer of Maersk.",
            ]),
            ("The Ask", [
                "Raising $13M Series A",
                "Use of funds: 50% product (ERP integrations, reporting modules) | 30% enterprise sales (EMEA focus) | 20% operations",
                "Two lead investors committed at $8M combined. Seeking strategic co-investor in ESG or enterprise software.",
                "Milestones: $12M ARR within 18 months, 30 enterprise customers, SAP partnership signed",
                "CSRD enforcement begins Jan 2025 for large EU companies — strongest pipeline quarter in company history",
            ]),
        ],
    },
]


# ---------------------------------------------------------------------------
# Generator
# ---------------------------------------------------------------------------

def generate_deck(deck: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        bg=deck["bg"],
        accent=deck["accent"],
        company=deck["company"],
        tagline=deck["tagline"],
        subtitle=deck["subtitle"],
    )

    for title, bullets in deck["slides"]:
        add_content_slide(prs, bg=deck["bg"], accent=deck["accent"], title=title, bullets=bullets)

    out_path = OUT_DIR / deck["filename"]
    prs.save(str(out_path))
    print(f"  Saved: {out_path}")


if __name__ == "__main__":
    print("Generating sample decks...")
    for deck in DECKS:
        generate_deck(deck)
    print("Done.")
