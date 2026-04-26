"""
Brand Intelligence — synchronous preprocessing step.

Extracts color palette, typography, and logo from the uploaded deck file.
Runs before the async pipeline on raw file bytes — adds no latency.

Returns a BrandProfile dict consumed by Skill 3 to prepend a Brand Guide
to the output and inject formatting hints into each slide comment.
"""

import base64
import colorsys
import io
import re
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


# ---------------------------------------------------------------------------
# Color utilities
# ---------------------------------------------------------------------------

def _rgb_to_hex(r: int, g: int, b: int) -> str:
    return f"#{r:02X}{g:02X}{b:02X}"


def _hex_to_hsl(hex_color: str) -> tuple[float, float, float]:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16) / 255, int(hex_color[2:4], 16) / 255, int(hex_color[4:6], 16) / 255
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    return h * 360, s * 100, l * 100


def _label_color(hex_color: str) -> str:
    """Classify a hex color by role based on HSL values."""
    h, s, l = _hex_to_hsl(hex_color)
    if l > 90:
        return "background"
    if l < 15:
        return "text-dark"
    if s > 55 and 25 < l < 75:
        return "accent"
    if l < 40:
        return "primary-dark"
    return "secondary"


def _deduplicate_colors(colors: list[str]) -> list[str]:
    """Remove near-duplicate colors (within ~15 lightness points)."""
    seen: list[str] = []
    for c in colors:
        h1, s1, l1 = _hex_to_hsl(c)
        is_dupe = any(
            abs(l1 - _hex_to_hsl(s)[2]) < 12 and abs(h1 - _hex_to_hsl(s)[0]) < 20
            for s in seen
        )
        if not is_dupe:
            seen.append(c)
    return seen


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------

def _extract_pptx(file_bytes: bytes) -> dict:
    prs = Presentation(io.BytesIO(file_bytes))
    raw_colors: list[str] = []
    fonts: list[str] = []
    logo_b64: str | None = None
    logo_ext: str = "png"

    def _collect_shape(shape, is_slide_one: bool):
        nonlocal logo_b64, logo_ext

        # Colors from shape fill
        try:
            fill = shape.fill
            if fill and fill.fore_color and fill.fore_color.rgb:
                rgb = fill.fore_color.rgb
                raw_colors.append(_rgb_to_hex((rgb >> 16) & 0xFF, (rgb >> 8) & 0xFF, rgb & 0xFF))
        except Exception:
            pass

        # Colors and fonts from text frames
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.name and run.font.name not in fonts:
                            fonts.append(run.font.name)
                    except Exception:
                        pass
                    try:
                        if run.font.color and run.font.color.rgb:
                            rgb = run.font.color.rgb
                            raw_colors.append(_rgb_to_hex((rgb >> 16) & 0xFF, (rgb >> 8) & 0xFF, rgb & 0xFF))
                    except Exception:
                        pass

        # Logo: first image on slide 1
        if is_slide_one and logo_b64 is None:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = shape.image.blob
                    ct = shape.image.content_type or "image/png"
                    logo_ext = ct.split("/")[-1].replace("jpeg", "jpg")
                    if len(img_bytes) <= 150_000:
                        logo_b64 = base64.b64encode(img_bytes).decode()
            except Exception:
                pass

    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            _collect_shape(shape, is_slide_one=(slide_num == 0))
        # Only scan first 3 slides for logo to avoid false positives
        if slide_num >= 2:
            logo_b64 = logo_b64  # freeze

    # Also check slide master for fonts
    try:
        for layout in prs.slide_master.slide_layouts[:3]:
            for shape in layout.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.name not in fonts:
                                fonts.append(run.font.name)
    except Exception:
        pass

    # Filter out near-white and pure-black (too common to be meaningful)
    meaningful_colors = [
        c for c in raw_colors
        if c not in ("#FFFFFF", "#000000", "#FEFEFE", "#010101")
    ]

    return {
        "source": "pptx",
        "colors": _deduplicate_colors(meaningful_colors)[:8],
        "fonts": fonts[:4],
        "logo_b64": logo_b64,
        "logo_ext": logo_ext,
    }


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------

def _extract_pdf(file_bytes: bytes) -> dict:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    fonts: list[str] = []
    raw_colors: list[str] = []
    logo_b64: str | None = None
    logo_ext: str = "png"

    for page_num, page in enumerate(doc):
        # Fonts
        for font_info in page.get_fonts():
            name = (font_info[3] or font_info[4] or "").split("+")[-1]  # strip subset prefix
            base_name = re.split(r"[-_,]", name)[0]  # Helvetica-Bold → Helvetica
            if base_name and base_name not in fonts and len(base_name) > 2:
                fonts.append(base_name)

        # Colors from text spans
        try:
            blocks = page.get_text("dict", flags=0)["blocks"]
            for block in blocks:
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            color_int = span.get("color", 0)
                            r = (color_int >> 16) & 0xFF
                            g = (color_int >> 8) & 0xFF
                            b = color_int & 0xFF
                            hex_c = _rgb_to_hex(r, g, b)
                            if hex_c not in ("#000000", "#FFFFFF") and hex_c not in raw_colors:
                                raw_colors.append(hex_c)
        except Exception:
            pass

        # Logo: first image on page 1
        if page_num == 0 and logo_b64 is None:
            for img in page.get_images():
                xref = img[0]
                try:
                    image_data = doc.extract_image(xref)
                    img_bytes = image_data["image"]
                    if len(img_bytes) <= 150_000:
                        logo_b64 = base64.b64encode(img_bytes).decode()
                        logo_ext = image_data.get("ext", "png")
                    break
                except Exception:
                    continue

        if page_num >= 2:
            break  # Only scan first 3 pages

    return {
        "source": "pdf",
        "colors": _deduplicate_colors(raw_colors)[:8],
        "fonts": fonts[:4],
        "logo_b64": logo_b64,
        "logo_ext": logo_ext,
    }


# ---------------------------------------------------------------------------
# Public interface
# ---------------------------------------------------------------------------

def extract(file_bytes: bytes, filename: str) -> dict:
    """Entry point. Returns a BrandProfile dict."""
    fn = filename.lower()
    try:
        if fn.endswith(".pptx"):
            return _extract_pptx(file_bytes)
        elif fn.endswith(".pdf"):
            return _extract_pdf(file_bytes)
    except Exception:
        pass
    return {"source": "unknown", "colors": [], "fonts": [], "logo_b64": None, "logo_ext": "png"}


def format_for_prompt(brand: dict) -> str:
    """
    Returns the brand section injected into Skill 3's generation prompt.
    Tells Claude what formatting instructions to include in the slide output.
    """
    if not brand.get("colors") and not brand.get("fonts"):
        return ""

    lines = ["BRAND & FORMATTING PROFILE (extracted from the uploaded deck):"]

    if brand.get("fonts"):
        lines.append(f"Typography: {', '.join(brand['fonts'])}")

    if brand.get("colors"):
        labelled = [(c, _label_color(c)) for c in brand["colors"]]
        color_str = "  |  ".join(f"{c} ({label})" for c, label in labelled)
        lines.append(f"Color palette: {color_str}")

    if brand.get("logo_b64"):
        lines.append("Logo: extracted and embedded in output")
    else:
        lines.append("Logo: not detected or too large to embed — founder should add manually")

    lines.append(
        "\nIn each slide's HTML comment, include a Formatting line with the relevant brand colors "
        "and font for that slide type. Example:\n"
        "<!-- Slide 1: Problem | Layout: Headline + 3 bullets | "
        "Formatting: bg #FFFFFF, heading #1A1A2E Montserrat Bold, accent #E94560 -->"
    )

    return "\n".join(lines)


def format_for_notebooklm_instructions(brand: dict) -> str:
    """
    Serialises the brand profile into a plain-text instruction string for
    NotebookLM's slide generation. NotebookLM won't apply these perfectly —
    it uses its own templates — but it will respect tone, color naming, and
    font preferences where its template allows.
    """
    if not brand.get("colors") and not brand.get("fonts"):
        return ""

    parts = ["Apply the following brand guidelines as closely as possible:"]

    if brand.get("fonts"):
        primary = brand["fonts"][0]
        rest = brand["fonts"][1:]
        parts.append(f"Primary font: {primary}" + (f". Secondary: {', '.join(rest)}" if rest else ""))

    if brand.get("colors"):
        labelled = [(c, _label_color(c)) for c in brand["colors"]]
        color_lines = ", ".join(f"{c} ({label})" for c, label in labelled)
        parts.append(f"Brand colors: {color_lines}")

    parts.append("Use these colors for slide backgrounds, headings, and accents where the template allows.")

    return " ".join(parts)


def format_brand_guide_markdown(brand: dict, company_name: str = "") -> str:
    """
    Returns the Brand Guide block prepended to the deck output.
    Includes logo if available.
    """
    if not brand.get("colors") and not brand.get("fonts") and not brand.get("logo_b64"):
        return ""

    lines = ["## Brand Guide", ""]
    lines.append("> Apply these brand settings when building slides from this narrative.\n")

    # Logo
    if brand.get("logo_b64"):
        ext = brand.get("logo_ext", "png").replace("jpg", "jpeg")
        mime = f"image/{ext}"
        lines.append(f'![{company_name or "Company"} Logo](data:{mime};base64,{brand["logo_b64"]})')
        lines.append("")

    # Typography
    if brand.get("fonts"):
        lines.append(f"**Typography:** {', '.join(brand['fonts'])}")

    # Color palette
    if brand.get("colors"):
        lines.append("\n**Color palette:**")
        for color in brand["colors"]:
            label = _label_color(color)
            lines.append(f"- `{color}` — {label}")

    lines.append("")
    lines.append("---")
    lines.append("")

    return "\n".join(lines)
