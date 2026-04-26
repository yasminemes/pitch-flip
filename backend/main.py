import asyncio
import csv
import io
import json
import os
from pathlib import Path

import fitz  # PyMuPDF
import httpx
from anthropic import Anthropic
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from duckduckgo_search import DDGS
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pptx import Presentation

from skills import brand_intelligence, skill1_company_intelligence, skill2_buyer_intelligence, skill3_synthesis

load_dotenv()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
PROMPTS = Path(__file__).parent / "prompts"
BULK_LIMIT = 30


# ---------------------------------------------------------------------------
# Deck parsing
# ---------------------------------------------------------------------------

def extract_pdf_text(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            pages.append(text)
    return "\n\n---\n\n".join(pages)


def extract_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    return "\n\n---\n\n".join(slides)


def parse_deck(file_bytes: bytes, filename: str) -> str:
    filename = filename.lower()
    if filename.endswith(".pdf"):
        return extract_pdf_text(file_bytes)
    elif filename.endswith(".pptx"):
        return extract_pptx_text(file_bytes)
    raise HTTPException(status_code=400, detail="Please upload a PDF or PPTX file.")


# ---------------------------------------------------------------------------
# Bulk upload helpers (unchanged)
# ---------------------------------------------------------------------------

def parse_csv(csv_bytes: bytes) -> list[dict]:
    text = csv_bytes.decode("utf-8-sig")
    reader = csv.DictReader(io.StringIO(text))
    prospects = []
    for row in reader:
        company = row.get("Customer Name", "").strip()
        if company:
            prospects.append({
                "company_name": company,
                "contact_name": row.get("Contact Name", "").strip(),
                "contact_title": row.get("Job Title", "").strip(),
            })
    return prospects[:BULK_LIMIT]


def parse_excel(excel_bytes: bytes) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes))
    ws = wb.active
    headers = [cell.value for cell in next(ws.iter_rows(min_row=1, max_row=1))]
    prospects = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        row_dict = dict(zip(headers, row))
        company = str(row_dict.get("Customer Name") or "").strip()
        if company:
            prospects.append({
                "company_name": company,
                "contact_name": str(row_dict.get("Contact Name") or "").strip(),
                "contact_title": str(row_dict.get("Job Title") or "").strip(),
            })
    return prospects[:BULK_LIMIT]


# ---------------------------------------------------------------------------
# Legacy enrich (still used by bulk)
# ---------------------------------------------------------------------------

async def enrich_company(company_name: str) -> str:
    def _fetch() -> str:
        parts = []
        try:
            with DDGS() as ddgs:
                overview = list(ddgs.text(f"{company_name} company what they do", max_results=3))
                if overview:
                    parts.append("COMPANY OVERVIEW:\n" + "\n".join(r["body"] for r in overview))

                news = list(ddgs.text(f"{company_name} strategy priorities challenges 2025 2026", max_results=2))
                if news:
                    parts.append("RECENT CONTEXT:\n" + "\n".join(r["body"] for r in news))

                website_results = list(ddgs.text(f"{company_name} official site", max_results=1))
                if website_results and website_results[0].get("href"):
                    try:
                        with httpx.Client(timeout=8, follow_redirects=True) as http:
                            resp = http.get(
                                website_results[0]["href"],
                                headers={"User-Agent": "Mozilla/5.0"},
                            )
                            soup = BeautifulSoup(resp.text, "html.parser")
                            for tag in soup(["script", "style", "nav", "footer", "header"]):
                                tag.decompose()
                            page_text = soup.get_text(separator=" ", strip=True)[:2000]
                            if page_text:
                                parts.append(f"FROM THEIR WEBSITE:\n{page_text}")
                    except Exception:
                        pass
        except Exception:
            pass
        return "\n\n---\n\n".join(parts) if parts else (
            f"No enrichment data found. Use your knowledge of {company_name} to personalise the narrative."
        )

    try:
        return await asyncio.wait_for(asyncio.to_thread(_fetch), timeout=20.0)
    except Exception:
        return f"No enrichment data found. Use your knowledge of {company_name} to personalise the narrative."


def build_bulk_prompt(
    deck_text: str,
    company_name: str = "",
    contact_name: str = "",
    contact_title: str = "",
    company_research: str = "",
) -> str:
    if company_name:
        contact_parts = []
        if contact_name:
            contact_parts.append(f"Contact name: {contact_name}")
        if contact_title:
            contact_parts.append(f"Job title: {contact_title}")
        contact_line = "\n".join(contact_parts) if contact_parts else "No contact details provided."
        return (
            (PROMPTS / "rewrite_tailored.txt")
            .read_text()
            .replace("{company_name}", company_name)
            .replace("{contact_line}", contact_line)
            .replace("{company_research}", company_research)
            .replace("{deck_content}", deck_text)
        )
    return (PROMPTS / "rewrite.txt").read_text().replace("{deck_content}", deck_text)


# ---------------------------------------------------------------------------
# SSE helper
# ---------------------------------------------------------------------------

def sse(payload: dict) -> str:
    return f"data: {json.dumps(payload)}\n\n"


# ---------------------------------------------------------------------------
# Single rewrite — three-skill pipeline
# ---------------------------------------------------------------------------

@app.post("/rewrite")
async def rewrite(
    file: UploadFile = File(...),
    company_name: str = Form(default=""),
    contact_title: str = Form(default=""),
    contact_department: str = Form(default=""),
):
    content = await file.read()
    deck_text = parse_deck(content, file.filename or "")

    if not deck_text.strip():
        raise HTTPException(
            status_code=422,
            detail="Could not extract readable text from this file. It may be image-only. Try exporting as a PDF with selectable text.",
        )

    company_name = company_name.strip()
    contact_title = contact_title.strip()
    contact_department = contact_department.strip()

    # Brand extraction is synchronous and fast — run before the async pipeline
    brand = brand_intelligence.extract(content, file.filename or "")

    async def generate():
        # --- Skill 1 + Skill 2 in parallel ---
        yield sse({"type": "status", "text": "Analysing your deck and researching your company..."})

        if company_name:
            yield sse({"type": "status", "text": f"Researching {company_name}..."})

        skill2_coro = (
            skill2_buyer_intelligence.run(company_name, contact_title, contact_department, client)
            if company_name
            else asyncio.sleep(0)
        )

        results = await asyncio.gather(
            skill1_company_intelligence.run(deck_text, client),
            skill2_coro,
            return_exceptions=True,
        )

        company_brief = results[0] if isinstance(results[0], str) else ""
        buyer_brief = results[1] if isinstance(results[1], str) else ""

        # --- Skill 3: build prompt and stream ---
        yield sse({"type": "status", "text": "Writing your deck..."})

        prompt = skill3_synthesis.build_prompt(
            deck_text=deck_text,
            company_brief=company_brief,
            buyer_brief=buyer_brief,
            company_name=company_name,
            contact_title=contact_title,
            contact_department=contact_department,
            brand_profile=brand_intelligence.format_for_prompt(brand),
        )

        # Prepend Brand Guide block before streaming the deck content
        brand_guide = brand_intelligence.format_brand_guide_markdown(brand, company_name)
        if brand_guide:
            yield sse({"type": "text", "text": brand_guide})

        # Bridge sync Anthropic stream into the async generator via a queue
        loop = asyncio.get_event_loop()
        queue: asyncio.Queue[str | None] = asyncio.Queue()

        def _sync_stream():
            with client.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=4096,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    loop.call_soon_threadsafe(queue.put_nowait, text)
            loop.call_soon_threadsafe(queue.put_nowait, None)

        stream_task = loop.run_in_executor(None, _sync_stream)

        while True:
            chunk = await queue.get()
            if chunk is None:
                break
            yield sse({"type": "text", "text": chunk})

        await stream_task

        # Send buyer brief so frontend can reveal it on demand
        if buyer_brief:
            yield sse({"type": "brief", "data": buyer_brief})

        yield "data: [DONE]\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ---------------------------------------------------------------------------
# Bulk rewrite (unchanged logic, uses legacy enrich)
# ---------------------------------------------------------------------------

@app.post("/rewrite-bulk")
async def rewrite_bulk(
    file: UploadFile = File(...),
    csv_file: UploadFile = File(...),
):
    deck_bytes = await file.read()
    csv_bytes = await csv_file.read()

    deck_text = parse_deck(deck_bytes, file.filename or "")
    if not deck_text.strip():
        raise HTTPException(
            status_code=422,
            detail="Could not extract readable text from the deck. It may be image-only.",
        )

    csv_filename = (csv_file.filename or "").lower()
    try:
        if csv_filename.endswith(".xlsx"):
            prospects = parse_excel(csv_bytes)
        else:
            prospects = parse_csv(csv_bytes)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not read file: {e}")

    if not prospects:
        raise HTTPException(
            status_code=422,
            detail="No valid prospects found in the CSV. Make sure it has a 'company_name' column.",
        )

    enriched = []
    for prospect in prospects:
        research = await enrich_company(prospect["company_name"])
        enriched.append((prospect, research))

    def generate():
        total = len(enriched)
        for i, (prospect, research) in enumerate(enriched):
            company_name = prospect["company_name"]

            yield f"data: {json.dumps({'type': 'start', 'company': company_name, 'index': i, 'total': total})}\n\n"

            prompt = build_bulk_prompt(
                deck_text,
                company_name=company_name,
                contact_name=prospect.get("contact_name", ""),
                contact_title=prospect.get("contact_title", ""),
                company_research=research,
            )

            with client.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=4096,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    yield f"data: {json.dumps({'type': 'text', 'index': i, 'text': text})}\n\n"

            yield f"data: {json.dumps({'type': 'done', 'index': i})}\n\n"

        yield "data: [DONE]\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.post("/slides")
async def generate_slides(
    deck_content: str = Form(...),
    file: UploadFile = File(default=None),
):
    """
    Accepts the generated deck markdown and optionally the original deck file.
    Creates a NotebookLM notebook, generates a PPTX, and streams SSE status +
    the final file. Brand data extracted from the original file is passed as
    instructions to NotebookLM as a partial branding workaround.
    """
    from notebooklm import NotebookLMClient
    from notebooklm._artifacts import SlideDeckFormat
    import base64
    import tempfile

    brand_instructions: str | None = None
    if file:
        try:
            file_bytes = await file.read()
            brand = brand_intelligence.extract(file_bytes, file.filename or "")
            brand_instructions = brand_intelligence.format_for_notebooklm_instructions(brand) or None
        except Exception:
            pass

    async def stream():
        nb_id = None
        tmp_path = None
        try:
            yield sse({"type": "status", "text": "Creating NotebookLM notebook..."})
            async with await NotebookLMClient.from_storage() as nb_client:
                nb = await nb_client.notebooks.create("PitchFlip Slides")
                nb_id = nb.id

                yield sse({"type": "status", "text": "Adding deck content to notebook..."})
                await nb_client.sources.add_text(nb_id, "Sales Deck", deck_content, wait=True)

                yield sse({"type": "status", "text": "Generating slides — this takes a few minutes..."})
                gen_status = await nb_client.artifacts.generate_slide_deck(
                    nb_id,
                    slide_format=SlideDeckFormat.PRESENTER_SLIDES,
                    instructions=brand_instructions,
                )

                try:
                    await nb_client.artifacts.wait_for_completion(nb_id, gen_status.task_id, timeout=600)
                except TimeoutError:
                    pass  # attempt download anyway — may have completed server-side

                yield sse({"type": "status", "text": "Downloading PPTX..."})
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                    tmp_path = tmp.name

                await nb_client.artifacts.download_slide_deck(nb_id, tmp_path, output_format="pptx")

                pptx_b64 = base64.b64encode(Path(tmp_path).read_bytes()).decode()
                yield sse({"type": "file", "data": pptx_b64, "filename": "sales-deck.pptx"})

        except Exception as e:
            yield sse({"type": "error", "text": str(e)})
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
            if nb_id:
                try:
                    async with await NotebookLMClient.from_storage() as nb_client:
                        await nb_client.notebooks.delete(nb_id)
                except Exception:
                    pass
            yield "data: [DONE]\n\n"

    return StreamingResponse(
        stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.get("/health")
def health():
    return {"status": "ok"}
